"""Create PPTX presentation files for download.

Tier: std-create (Tier 1) — single-shot file generation.

**std-create**: build a PowerPoint from a list of markdown slides in
one call. Auto-loaded for every `agentino.tools.std` agent.

For multi-step deck work — open existing decks, replace placeholders,
apply slide templates — the convention is **Tier 2 / `skill-pptx`**:
a per-tenant skill bundle, opt-in via `skills:` in that tenant's
agents.yml.

Tool name in agents.yml: `create_presentation` (back-compat).
"""

import io
import re
from datetime import datetime

from agentino.core.tool import tool

from .storage import get_default_store


@tool
async def create_presentation(title: str, slides: list[str], filename: str = "") -> str:
    """Create a PowerPoint (.pptx) presentation and attach it for download. Each slide is a markdown string. Use # for slide titles, ## for subtitles, - for bullet points.

    Args:
        title: Presentation title (first slide)
        slides: List of slide contents in markdown format. Each string becomes one slide.
        filename: Optional output filename (without extension). Defaults to title-based name.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    tf = slide.shapes.title.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0x1F, 0x29, 0x37)

    if slide.placeholders[1]:
        sub = slide.placeholders[1].text_frame
        sub.paragraphs[0].text = datetime.now().strftime("%B %d, %Y")
        sub.paragraphs[0].font.size = Pt(18)
        sub.paragraphs[0].font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    # Content slides
    for slide_md in slides:
        lines = slide_md.strip().split("\n")
        slide_title = ""
        items = []

        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue
            if stripped.startswith("# "):
                slide_title = stripped[2:]
            elif stripped.startswith("## "):
                slide_title = stripped[3:]
            elif stripped.startswith("- ") or stripped.startswith("* "):
                items.append(("bullet", stripped[2:]))
            elif re.match(r"^\d+\.\s+", stripped):
                items.append(("number", re.sub(r"^\d+\.\s+", "", stripped)))
            else:
                items.append(("text", stripped))

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide_title and slide.shapes.title:
            tf = slide.shapes.title.text_frame
            tf.paragraphs[0].text = slide_title
            tf.paragraphs[0].font.size = Pt(32)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor(0x1F, 0x29, 0x37)

        if items and len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, (kind, text) in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                # Parse bold markers
                parts = re.split(r"(\*\*.*?\*\*)", text)
                for part in parts:
                    if part.startswith("**") and part.endswith("**"):
                        run = p.add_run()
                        run.text = part[2:-2]
                        run.font.bold = True
                    else:
                        run = p.add_run()
                        run.text = part

                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
                p.space_after = Pt(8)

                if kind == "bullet":
                    p.level = 0
                elif kind == "number":
                    p.level = 0

    if not filename:
        safe = title.lower().replace(" ", "-").replace("/", "-")[:40]
        filename = safe
    now_short = datetime.now().strftime("%Y-%m-%d")
    final_filename = f"{filename}-{now_short}.pptx"

    buf = io.BytesIO()
    prs.save(buf)
    content_bytes = buf.getvalue()

    # Indexable preview: title + slide markdown stripped of code fences.
    extracted = title + "\n\n" + "\n\n---\n\n".join(slides)
    extracted = re.sub(r"```[a-z]*\n.*?```", "", extracted, flags=re.DOTALL)[:50_000]

    try:
        stored = get_default_store().save(
            content_bytes=content_bytes,
            filename=final_filename,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            title=title,
            extracted_text=extracted,
        )
    except Exception as e:
        return f"Presentation generation failed: {e}"

    return (
        f"Presentation created ({len(slides)} slides). Download: [{stored.filename}]({stored.url})"
    )
