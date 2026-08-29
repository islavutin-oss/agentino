"""Read uploaded files — auto-detect format, extract content.

Reads through `protocols.FileStorage` (Supabase Storage / local).
Bytes are materialized to a temp file so the existing format-specific
readers (which take a Path) work unchanged.

History: previously walked the local workspace-files directory
directly, breaking tenant isolation and pinning the tool to local
disk. Migrated 2026-05-02.
"""

from __future__ import annotations

import json
import tempfile
from pathlib import Path

from agentino.core.tool import tool

# ── Format-specific readers (unchanged from pre-migration) ─────────────


def _read_csv(path: Path) -> str:
    import csv

    with open(path, newline="", encoding="utf-8-sig") as f:
        reader = csv.reader(f)
        rows = list(reader)
    if not rows:
        return "Empty CSV file."
    columns = rows[0]
    data = rows[1:]
    table = json.dumps({"columns": columns, "rows": data[:200]})
    summary = f"CSV: {len(data)} rows, {len(columns)} columns"
    if len(data) > 200:
        summary += f" (showing first 200 of {len(data)})"
    return f"{summary}\n\n```datatable\n{table}\n```"


def _read_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    results = []
    for ws in wb.worksheets:
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        columns = [str(c) if c is not None else "" for c in rows[0]]
        data = [[str(c) if c is not None else "" for c in row] for row in rows[1:201]]
        table = json.dumps({"title": ws.title, "columns": columns, "rows": data})
        summary = f"Sheet '{ws.title}': {len(rows) - 1} rows, {len(columns)} columns"
        if len(rows) - 1 > 200:
            summary += " (showing first 200)"
        results.append(f"{summary}\n\n```datatable\n{table}\n```")
    wb.close()
    return "\n\n".join(results) if results else "Empty spreadsheet."


def _read_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        headers = [cell.text for cell in table.rows[0].cells]
        rows = []
        for row in table.rows[1:]:
            rows.append([cell.text for cell in row.cells])
        if headers:
            table_json = json.dumps({"columns": headers, "rows": rows})
            parts.append(f"```datatable\n{table_json}\n```")
    return "\n".join(parts) if parts else "Empty document."


def _read_pdf(path: Path) -> str:
    try:
        import fitz

        doc = fitz.open(str(path))
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        return text.strip() or "Empty PDF (no extractable text)."
    except ImportError:
        pass
    try:
        import pdfplumber

        with pdfplumber.open(str(path)) as pdf:
            return (
                "\n".join(page.extract_text() or "" for page in pdf.pages).strip() or "Empty PDF."
            )
    except ImportError:
        return "[PDF reading requires PyMuPDF or pdfplumber — not installed]"


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        texts.append(" | ".join(cell.text for cell in row.cells))
            if texts:
                parts.append(f"Slide {i}:\n" + "\n".join(texts))
        return "\n\n".join(parts) if parts else "Empty presentation."
    except ImportError:
        return "[PPTX reading requires python-pptx — not installed]"


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


READERS = {
    ".csv": _read_csv,
    ".xlsx": _read_xlsx,
    ".xls": _read_xlsx,
    ".docx": _read_docx,
    ".pptx": _read_pptx,
    ".pdf": _read_pdf,
    ".txt": _read_text,
    ".md": _read_text,
    ".json": _read_text,
    ".html": _read_text,
    ".xml": _read_text,
    ".yaml": _read_text,
    ".yml": _read_text,
}


def _materialize_to_path(file_name: str, tenant_id: str) -> tuple[Path, str] | None:
    """Resolve `file_name` (file_id OR original_name) → temp Path + suffix.

    Returns (path, suffix) or None if the file isn't in storage. The
    caller should treat the path as read-only (tempfile-managed).
    """
    from agentino.tools.std._file_storage import get_file_storage

    storage = get_file_storage()
    content: bytes | None = None
    original_name = file_name

    # 1. Try file_id match
    try:
        content = storage.get(tenant_id, file_name)
    except (FileNotFoundError, ValueError):
        # 2. Walk listing for original_name match (LLM only sees that)
        try:
            for meta in storage.list(tenant_id):
                if (
                    meta.original_name == file_name
                    or meta.file_id == file_name
                    or file_name.lower() in (meta.original_name or "").lower()
                ):
                    content = storage.get(tenant_id, meta.file_id)
                    original_name = meta.original_name or file_name
                    break
        except Exception:
            pass

    if content is None:
        return None

    suffix = Path(original_name).suffix.lower()
    tmp = tempfile.NamedTemporaryFile(
        suffix=suffix or ".bin",
        prefix="read_",
        delete=False,
    )
    tmp.write(content)
    tmp.close()
    return Path(tmp.name), suffix


@tool(is_read_only=True)
async def read_file(file_name: str) -> str:
    """Read an uploaded file and return its content. Supports CSV, XLSX,
    DOCX, PDF, TXT, JSON, and more. For spreadsheets, returns data in
    table format.

    Args:
        file_name: filename or file_id (as returned by the upload
                   endpoint). Tenant scope is read from agentino.core.context.
    """
    from agentino.core.context import get_context

    tenant_id = get_context("tenant_id") or "default"

    resolved = _materialize_to_path(file_name, tenant_id)
    if resolved is None:
        return (
            f"File '{file_name}' not found in workspace files. "
            "Make sure the file was uploaded first."
        )

    path, suffix = resolved
    reader = READERS.get(suffix)
    if not reader:
        return f"Unsupported file format: {suffix}. Supported: {', '.join(READERS.keys())}"

    try:
        return reader(path)
    except Exception as e:
        return f"Error reading file: {e}"
